from pptx import Presentation
from core.schemas import DeckSpec, TemplateProfile

def render_pptx(template_path: str, deck_spec: DeckSpec, output_path: str, profile: TemplateProfile):
    """Renders the python-pptx presentation and saves it to output_path"""
    prs = Presentation(template_path)
    
    # Build dictionary of shape mapping to make applying fields easy
    # layout_id -> { field_key -> idx }
    layout_map = {}
    for layout in profile.layouts:
        layout_map[layout.layout_id] = {p.key: p.idx for p in layout.placeholders}
    
    for slide_spec in deck_spec.slides:
        layout_id = slide_spec.layout_id
        if layout_id < 0 or layout_id >= len(prs.slide_layouts):
            raise ValueError(f"Invalid layout_id: {layout_id}")
            
        layout = prs.slide_layouts[layout_id]
        slide = prs.slides.add_slide(layout)
        
        ph_map = layout_map.get(layout_id, {})
        
        for field in slide_spec.fields:
            field_key = field.key
            field_val = field.value
            if field_key not in ph_map:
                # MissingPlaceholderError analogue
                print(f"Warning: Placeholder '{field_key}' not found in layout {layout_id}. Skipping.")
                continue
                
            idx = ph_map[field_key]
            try:
                shape = slide.placeholders[idx]
            except KeyError:
                print(f"Warning: Shape index {idx} not found in layout {layout_id}. Skipping.")
                continue
            
            # Apply content
            if isinstance(field_val, list):
                # Body bullets
                text_frame = shape.text_frame
                text_frame.clear()  # removes all paragraphs
                for i, bullet_text in enumerate(field_val):
                    p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                    p.text = bullet_text
                    p.level = 0
            else:
                # Standard text
                shape.text_frame.text = str(field_val)
                
        # Notes
        if slide_spec.notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_spec.notes
            
    prs.save(output_path)

import os
import subprocess
from pdf2image import convert_from_path

def export_to_thumbnails(pptx_path: str, output_dir: str):
    """
    Converts a PPTX file to a series of JPEG thumbnails using LibreOffice and pdf2image.
    """
    os.makedirs(output_dir, exist_ok=True)
    
    # 1. Convert PPTX to PDF using LibreOffice headless
    pdf_path = pptx_path.replace(".pptx", ".pdf")
    try:
        # Mac OS path for LibreOffice or using 'soffice' if linked
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf", "--outdir", os.path.dirname(os.path.abspath(pptx_path)), pptx_path],
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE
        )
    except subprocess.CalledProcessError as e:
        print(f"LibreOffice conversion failed: {e.stderr.decode()}")
        return []
    except FileNotFoundError:
        print("LibreOffice (soffice) not found on path.")
        return []

    # 2. Convert PDF to Images
    if not os.path.exists(pdf_path):
        print("PDF was not created.")
        return []
        
    images = convert_from_path(pdf_path)
    output_files = []
    
    for i, image in enumerate(images):
        out_path = os.path.join(output_dir, f"slide_{i+1}.jpg")
        image.save(out_path, "JPEG")
        output_files.append(out_path)
        
    return output_files

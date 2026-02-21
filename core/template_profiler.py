from pptx import Presentation
from core.schemas import TemplateProfile, LayoutInfo, PlaceholderInfo
import os

def profile_template(template_path: str, template_name: str) -> TemplateProfile:
    """Parses a uploaded template and extracts the available layouts and their placeholders."""
    prs = Presentation(template_path)
    layouts = []
    
    for idx, layout in enumerate(prs.slide_layouts):
        placeholders = []
        for shape in layout.placeholders:
            # Attempt to derive a sensible key based on the placeholder name or type
            name = shape.name.lower()
            if "title" in name:
                key = "title"
            elif "subtitle" in name:
                key = "subtitle"
            elif "body" in name or "content" in name or "text" in name:
                key = "body"
            elif "footer" in name:
                key = "footer"
            elif "date" in name:
                key = "date"
            else:
                key = f"ph_{shape.placeholder_format.idx}"

            ph = PlaceholderInfo(
                key=key,
                type=shape.placeholder_format.type.__name__ if hasattr(shape.placeholder_format.type, '__name__') else str(shape.placeholder_format.type),
                idx=shape.placeholder_format.idx
            )
            placeholders.append(ph)
        
        # Deduplicate keys if necessary (e.g. multiple body placeholders)
        key_counts = {}
        for ph in placeholders:
            if ph.key in key_counts:
                key_counts[ph.key] += 1
                ph.key = f"{ph.key}_{key_counts[ph.key]}"
            else:
                key_counts[ph.key] = 1

        li = LayoutInfo(
            layout_id=idx,
            layout_name=layout.name,
            placeholders=placeholders
        )
        layouts.append(li)
    
    # By default, allow all layouts
    allowed_ids = [l.layout_id for l in layouts]
    
    return TemplateProfile(
        template_name=template_name,
        layouts=layouts,
        allowed_layout_ids=allowed_ids
    )

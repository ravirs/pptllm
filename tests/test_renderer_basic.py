from core.renderer import render_pptx
from core.schemas import DeckSpec, TemplateProfile, LayoutInfo, PlaceholderInfo
import os
from pptx import Presentation

def test_renderer_creates_file(tmp_path):
    # Create a blank presentation to use as a "template"
    dummy_template = tmp_path / "dummy.pptx"
    prs = Presentation()
    prs.save(dummy_template)
    
    # Needs to match the dummy presentation's actual default layouts
    # prs.slide_layouts[0] usually has a title (0) and subtitle (1)
    
    profile = TemplateProfile(
        template_name="dummy.pptx",
        layouts=[
            LayoutInfo(
                layout_id=0,
                layout_name="Title Slide",
                placeholders=[
                    PlaceholderInfo(key="title", type="TITLE", idx=0),
                    PlaceholderInfo(key="subtitle", type="SUBTITLE", idx=1)
                ]
            )
        ],
        allowed_layout_ids=[0]
    )
    
    deck = DeckSpec(
        deck_title="Test",
        slides=[
            {
                "slide_id": "s1",
                "layout_id": 0,
                "fields": [
                    {"key": "title", "value": "Hello"},
                    {"key": "subtitle", "value": "World"}
                ]
            }
        ]
    )
    
    output_path = tmp_path / "output.pptx"
    
    render_pptx(str(dummy_template), deck, str(output_path), profile)
    
    assert output_path.exists()
    
    # verify
    out_prs = Presentation(output_path)
    assert len(out_prs.slides) == 1
    slide = out_prs.slides[0]
    # In a default template, shapes[0] is title, shapes[1] is subtitle
    assert slide.shapes[0].text == "Hello"    
